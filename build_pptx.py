#!/usr/bin/env python3
"""
Generator Berkas Microsoft PowerPoint (.PPTX)
Makalah: Ahli Waris Laki-laki dan Perempuan
Mata Kuliah: Hukum Kewarisan Islam
Kelompok 5 - UIN Siber Syekh Nurjati Cirebon 2026
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Theme Colors: Emerald Forest & Imperial Gold
    BG_DARK = RGBColor(6, 35, 28)      # #06231C
    CARD_BG = RGBColor(11, 56, 44)     # #0B382C
    GOLD = RGBColor(230, 202, 101)     # #E6CA65
    EMERALD = RGBColor(16, 185, 129)   # #10B981
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(203, 213, 225)
    MUTED_GRAY = RGBColor(148, 163, 184)
    RED_ACCENT = RGBColor(239, 68, 68)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category="HUKUM KEWARISAN ISLAM • KELOMPOK 5"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        
        # Category Tag
        p0 = tf.paragraphs[0]
        p0.text = category.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = GOLD
        p0.font.name = "Arial"

        # Title
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = WHITE
        p1.font.name = "Arial"

    # ==========================================
    # SLIDE 1: COVER
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Decorative Border
    border = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(6.5))
    border.fill.solid()
    border.fill.fore_color.rgb = CARD_BG
    border.line.color.rgb = GOLD
    border.line.width = Pt(2)

    # Title content
    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(0.9), Inches(10.9), Inches(3.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "بِسْمِ اللَّهِ الرَّحْمَٰنِ الرَّحِيمِ"
    p.font.size = Pt(26)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "AHLI WARIS LAKI-LAKI & PEREMPUAN"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "Kajian Komprehensif Kedudukan, Klasifikasi, Dasar Hukum, dan Pembagian Warisan dalam Syariat Islam & KHI"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Info Box
    info_rect = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.8), Inches(4.3), Inches(9.7), Inches(2.2))
    info_rect.fill.solid()
    info_rect.fill.fore_color.rgb = RGBColor(4, 20, 15)
    info_rect.line.color.rgb = GOLD
    info_rect.line.width = Pt(1)

    tb_info = s1.shapes.add_textbox(Inches(2.0), Inches(4.4), Inches(9.3), Inches(2.0))
    tfi = tb_info.text_frame
    tfi.word_wrap = True

    p = tfi.paragraphs[0]
    p.text = "DOSEN PENGAMPU: Prof. Dr. H. Kosim, M.Ag"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p = tfi.add_paragraph()
    p.text = "Disusun Oleh Kelompok 5:\nFatiaful Alzahra (2530311087) • Sufyan Tsaury (2530311086) • Wahdan Hamdun (2530311003)"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tfi.add_paragraph()
    p.text = "PRODI HUKUM KELUARGA • FAKULTAS SYARI'AH • UIN SIBER SYEKH NURJATI CIREBON (2026)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 2: LATAR BELAKANG & URGENSI
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Latar Belakang & Urgensi Kewarisan Islam", "BAB I • PENDAHULUAN")

    cards_data_s2 = [
        ("1. Realitas & Sengketa", "Dalam kehidupan masyarakat, persoalan kewarisan sangat rentan memicu konflik keluarga akibat minimnya pemahaman tentang siapa yang berhak dan siapa yang terhalang (mahjub).", GOLD),
        ("2. Tiga Rukun Pokok Faraidh", "• Al-Muwarrits: Pewaris yang wafat.\n• Al-Waris: Ahli waris yang masih hidup saat pewaris wafat.\n• Al-Mauruts / Tirkah: Harta bersih setelah utang, zakat, & jenazah.", EMERALD),
        ("3. Kepastian Hukum Syariat", "Hukum waris Islam dirancang langsung oleh Allah SWT secara pasti (qath'iy) untuk menegakkan keadilan distributif bagi laki-laki maupun perempuan secara proporsional.", GOLD)
    ]

    for i, (title, content, border_color) in enumerate(cards_data_s2):
        x = Inches(0.8 + i * 4.0)
        c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(5.0))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = border_color
        c.line.width = Pt(1.5)

        tb = s2.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = border_color

        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(12)
        p2.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 3: RUMUSAN MASALAH & TUJUAN
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "Rumusan Masalah & Tujuan Penulisan", "FOKUS KAJIAN MAKALAH")

    # Left: Rumusan Masalah
    c_left = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = GOLD

    tb_l = s3.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "4 RUMUSAN MASALAH UTAMA"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_l.add_paragraph()
    p.text = "1. Apa pengertian ahli waris dalam hukum kewarisan Islam?\n\n2. Siapa sajakah yang termasuk ahli waris laki-laki?\n\n3. Siapa sajakah yang termasuk ahli waris perempuan?\n\n4. Bagaimana kedudukan hak ahli waris laki-laki dan perempuan (Ashhabul Furudh, 'Ashabah, dan Hijab)?"
    p.font.size = Pt(12.5)
    p.font.color.rgb = WHITE

    # Right: Tujuan Penulisan
    c_right = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = EMERALD

    tb_r = s3.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "CAPAIAN & TUJUAN PENULISAN"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_r.add_paragraph()
    p.text = "1. Memahami konsep yuridis & normatif kewarisan Islam secara komprehensif.\n\n2. Mengidentifikasi dan memetakan 15 golongan ahli waris laki-laki.\n\n3. Mengidentifikasi dan memetakan 10 golongan ahli waris perempuan.\n\n4. Mengetahui implikasi praktis dan penyelesaian sengketa waris menurut Fikih & KHI."
    p.font.size = Pt(12.5)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 4: PENGERTIAN & SEBAB MEWARISI
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Pengertian, Syarat, dan Sebab-Sebab Mewarisi", "BAB II • PEMBAHASAN DASAR")

    c1 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = GOLD

    tb_c1 = s4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "HAKIKAT AHLI WARIS & SYARAT"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_c1.add_paragraph()
    p.text = "Ahli Waris (Al-Waris) adalah orang yang mempunyai hak legal syariah untuk menerima harta peninggalan dari pewaris (al-muwarrits) karena adanya sebab yang diakui syariat.\n\n3 Syarat Sah Mewarisi:\n1. Kematian Pewaris: Meninggal hakiki atau putusan pengadilan.\n2. Hidupnya Ahli Waris: Masih bernyawa saat pewaris wafat.\n3. Tidak Ada Penghalang (Mani'): Bersih dari pembunuhan, beda agama, dll."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    c2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = EMERALD

    tb_c2 = s4.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "DUA SEBAB POKOK KEWARISAN"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_c2.add_paragraph()
    p.text = "1. Hubungan Nasabiyah (Pertalian Darah):\n• Ushul (Garis Atas): Ayah, Ibu, Kakek, Nenek\n• Furu' (Garis Bawah): Anak, Cucu\n• Hawasyi (Garis Samping): Saudara, Paman, Keponakan\n\n2. Hubungan Sababiyah (Peristiwa Hukum Sah):\n• Perkawinan yang Sah (Suami & Istri saling mewarisi)\n• Al-Wala' / Maula (Pembebasan budak - Fikih klasik)"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 5: DASAR HUKUM AL-QUR'AN
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "Dasar Hukum: Landasan Norma Ayat-Ayat Al-Qur'an", "SUMBER HUKUM TERTINGGI")

    ayats = [
        ("QS. An-Nisa Ayat 7", "Prinsip Kesetaraan Hak Mewarisi", "Bagi laki-laki ada hak bagian dari harta peninggalan orang tua & kerabat, dan bagi perempuan ada hak bagian (pula)... baik sedikit atau banyak menurut bagian yang telah ditetapkan.", GOLD),
        ("QS. An-Nisa Ayat 11", "Kewarisan Anak & Orang Tua", "Allah mensyariatkan bagimu tentang (warisan) anak-anakmu, yaitu bagian seorang anak laki-laki sama dengan bagian dua orang anak perempuan (2:1)...", EMERALD),
        ("QS. An-Nisa Ayat 12", "Hak Waris Suami & Istri", "Dan bagianmu (suami) 1/2 jika tak ada anak; 1/4 jika ada anak. Bagian istri 1/4 jika tak ada anak; 1/8 jika ada anak. Serta ketentuan saudara seibu.", GOLD),
        ("QS. An-Nisa Ayat 176", "Kewarisan Kalalah", "Mereka meminta fatwa kepadamu (tentang kalalah). Mengatur ketentuan pembagian warisan bagi saudara kandung dan seayah saat pewaris tak punya keturunan & orang tua.", EMERALD),
    ]

    for i, (title, sub, text, col) in enumerate(ayats):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + col_idx * 6.0)
        y = Inches(1.7 + row * 2.6)

        c = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.3))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col

        tb = s5.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{title} • {sub}"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = f'"{text}"'
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 6: HADIS & KHI
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "Dasar Hukum: Hadis Nabawi & Kompilasi Hukum Islam (KHI)", "YURISPRUDENSI DAN FIKIH")

    c1 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = GOLD

    tb_h = s6.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_h = tb_h.text_frame
    tf_h.word_wrap = True
    p = tf_h.paragraphs[0]
    p.text = "HADIS NABI SAW (BUKHARI & MUSLIM)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_h.add_paragraph()
    p.text = 'أَلْحِقُوا الْفَرَائِضَ بِأَهْلِهَا ، فَمَا بَقِيَ فَهُوَ لِأَوْلَى رَجُلٍ ذَكَرٍ'
    p.font.size = Pt(16)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.RIGHT

    p = tf_h.add_paragraph()
    p.text = '"Bagikanlah harta warisan kepada orang yang berhak menerimanya (Ashhabul Furudh), dan apa yang tersisa maka berikanlah kepada kerabat laki-laki yang paling dekat (\'Ashabah)."\n\nPrinsip:\n1. Hak pasti Ashhabul Furudh wajib dipenuhi terlebih dahulu.\n2. Sisanya diserahkan kepada \'Ashabah terdekat.'
    p.font.size = Pt(11.5)
    p.font.color.rgb = WHITE

    c2 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = EMERALD

    tb_khi = s6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_khi = tb_khi.text_frame
    tf_khi.word_wrap = True
    p = tf_khi.paragraphs[0]
    p.text = "KOMPILASI HUKUM ISLAM (KHI BUKU II)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_khi.add_paragraph()
    p.text = "• Pasal 171 Huruf c KHI:\nAhli waris adalah orang yang saat pewaris meninggal mempunyai hubungan darah atau perkawinan, beragama Islam, dan tidak terhalang hukum.\n\n• Pasal 174 KHI:\n1. Hubungan Darah:\n   - Pria: Ayah, anak laki-laki, saudara, paman, kakek.\n   - Wanita: Ibu, anak perempuan, saudara perempuan, nenek.\n2. Hubungan Perkawinan:\n   - Duda (Suami) & Janda (Istri)."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 7: 15 AHLI WARIS LAKI-LAKI
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "15 Golongan Ahli Waris Laki-Laki (Rijal)", "KLASIFIKASI AHLI WARIS")

    c_box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = CARD_BG
    c_box.line.color.rgb = GOLD

    tb_w15 = s7.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.8))
    tf_w15 = tb_w15.text_frame
    tf_w15.word_wrap = True

    p = tf_w15.paragraphs[0]
    p.text = "DAFTAR 15 AHLI WARIS LAKI-LAKI DAN KEDUDUKANNYA"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_w15.add_paragraph()
    p.text = "1. Anak Laki-laki (Ibn) - 'Ashabah terkuat       | 9. Keponakan Lk (dari Sdr Seayah)\n2. Cucu Laki-laki dari anak laki-laki          | 10. Paman Sekandung ('Amm Syaqiq)\n3. Ayah (Ab) - Fardh / 'Ashabah                | 11. Paman Seayah ('Amm li Ab)\n4. Kakek Shahih (Ayah dari Ayah)               | 12. Sepupu Lk dari Paman Sekandung\n5. Saudara Laki-laki Sekandung                 | 13. Sepupu Lk dari Paman Seayah\n6. Saudara Laki-laki Seayah                    | 14. Suami (Zawj) - 1/2 atau 1/4\n7. Saudara Laki-laki Seibu (1/6 atau 1/3)      | 15. Mu'tiq (Laki-laki pembebas budak)\n8. Keponakan Lk (dari Sdr Kandung)"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    p = tf_w15.add_paragraph()
    p.text = "★ KAIDAH PENTING: Jika ke-15 ahli waris laki-laki ada bersamaan, hanya 3 orang yang pasti mewarisi: Anak Laki-laki, Ayah, dan Suami."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # ==========================================
    # SLIDE 8: 10 AHLI WARIS PEREMPUAN
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "10 Golongan Ahli Waris Perempuan (Nisa)", "KLASIFIKASI AHLI WARIS")

    c_box8 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2))
    c_box8.fill.solid()
    c_box8.fill.fore_color.rgb = CARD_BG
    c_box8.line.color.rgb = EMERALD

    tb_w10 = s8.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.8))
    tf_w10 = tb_w10.text_frame
    tf_w10.word_wrap = True

    p = tf_w10.paragraphs[0]
    p.text = "DAFTAR 10 AHLI WARIS PEREMPUAN DAN HAKNYA"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_w10.add_paragraph()
    p.text = "1. Anak Perempuan (Bint) - 1/2, 2/3, atau 'Ashabah bi Ghairih bersama anak laki-laki\n2. Cucu Perempuan (Bint al-Ibn) - Dari jalur anak laki-laki (1/2, 2/3, atau pelengkap 1/6)\n3. Ibu (Umm) - 1/3 (tanpa anak/saudara) atau 1/6 (ada anak/beberapa saudara) [PASTI DAPAT]\n4. Nenek dari Pihak Ibu (Jaddah li Umm) - 1/6 jika tidak ada ibu\n5. Nenek dari Pihak Ayah (Jaddah li Ab) - 1/6 jika tidak ada ibu dan ayah\n6. Saudara Perempuan Sekandung (Ukht Syaqiqah) - 1/2, 2/3, atau 'Ashabah\n7. Saudara Perempuan Seayah (Ukht li Ab) - 1/2, 2/3, 1/6, atau 'Ashabah\n8. Saudara Perempuan Seibu (Ukht li Umm) - 1/6 (sendiri) atau 1/3 (bersama)\n9. Istri (Zawjah) - 1/4 (tanpa anak) atau 1/8 (ada anak)\n10. Mu'tiqah (Wanita yang memerdekakan budak - Fikih klasik)"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    p = tf_w10.add_paragraph()
    p.text = "★ KAIDAH PENTING: Jika ke-10 ahli waris perempuan ada bersamaan, 5 orang yang pasti mewarisi: Istri, Anak Perempuan, Ibu, Cucu Perempuan, dan Saudara Perempuan Sekandung."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # ==========================================
    # SLIDE 9: ASHHABUL FURUDH VS 'ASHABAH
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "Sistem Hak: Ashhab al-Furudh & 'Ashabah", "MEKANISME DISTRIBUSI WARIS")

    c1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = GOLD

    tb_f = s9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    p = tf_f.paragraphs[0]
    p.text = "ASHHAB AL-FURUDH (6 Bagian Pasti)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_f.add_paragraph()
    p.text = "• 1/2 : Suami (tanpa anak), 1 Anak Pr, 1 Cucu Pr, 1 Sdr Pr Kandung/Seayah.\n• 1/4 : Suami (ada anak) atau Istri (tanpa anak).\n• 1/8 : Istri (jika pewaris memiliki keturunan).\n• 2/3 : 2 atau lebih Anak Pr, Cucu Pr, atau Sdr Pr Kandung/Seayah.\n• 1/3 : Ibu (tanpa anak & saudara), atau >=2 Sdr Seibu.\n• 1/6 : Ayah/Ibu (ada anak), Kakek/Nenek, 1 Sdr Seibu, Cucu Pr pelengkap 2/3."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    c2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = EMERALD

    tb_a = s9.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    p = tf_a.paragraphs[0]
    p.text = "AL-'ASHABAH (Penerima Sisa Harta)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_a.add_paragraph()
    p.text = "1. 'Ashabah bi Nafsihi:\nKerabat laki-laki garis nasab langsung tanpa perantara wanita (Anak lk, ayah, paman, saudara sekandung/seayah).\n\n2. 'Ashabah bi Ghairihi:\nWanita yang ditarik menjadi penerima sisa oleh saudara laki-lakinya (Anak pr ditarik anak lk dengan rasio 2:1).\n\n3. 'Ashabah ma'a Ghairihi:\nSaudara perempuan kandung/seayah bersama dengan anak/cucu perempuan."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 10: MANI' & HIJAB
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "Penghalang Hak Waris (Mani') & Konsep Hijab", "BATASAN DAN IMPEDIMEN")

    c1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = RED_ACCENT

    tb_m = s10.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    p = tf_m.paragraphs[0]
    p.text = "3 FAKTOR PENGHALANG (MANI' AL-IRTS)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT

    p = tf_m.add_paragraph()
    p.text = "1. Pembunuhan (Al-Qatlu):\nAhli waris yang membunuh pewaris tidak berhak mendapatkan warisan apa pun (HR. Tirmidzi & Ahmad).\n\n2. Perbedaan Agama (Ikhtilaf ad-Din):\nSeorang muslim tidak mewarisi orang non-muslim, dan sebaliknya (HR. Bukhari & Muslim).\n\n3. Perbudakan (Ar-Riqq):\nBudak tidak cakap hukum memiliki harta (pembahasan historis fikih klasik)."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    c2 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = EMERALD

    tb_hij = s10.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_hij = tb_hij.text_frame
    tf_hij.word_wrap = True
    p = tf_hij.paragraphs[0]
    p.text = "KONSEP AL-HIJAB (DINDING PENGHALANG)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_hij.add_paragraph()
    p.text = "1. Hijab Nuqshan (Pengurangan Porsi):\nPengurangan bagian karena kehadiran ahli waris lain:\n• Suami: 1/2 -> 1/4 (karena ada anak)\n• Istri: 1/4 -> 1/8 (karena ada anak)\n• Ibu: 1/3 -> 1/6 (karena ada anak/beberapa saudara)\n\n2. Hijab Hirman (Gugur Total):\nTertutupnya hak waris karena adanya kerabat yang lebih dekat:\n• Kakek terhalang oleh Ayah\n• Cucu laki-laki terhalang oleh Anak Laki-laki\n• Saudara terhalang oleh Ayah atau Anak Laki-laki."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 11: KOMPARASI & KEADILAN
    # ==========================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)
    add_header(s11, "Komparasi & Falsafah Keadilan Proporsional", "FILOSOFI HUKUM ISLAM")

    c1 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = GOLD

    tb_tab = s11.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_tab = tb_tab.text_frame
    tf_tab.word_wrap = True
    p = tf_tab.paragraphs[0]
    p.text = "TABEL KOMPARASI HAK MEWARISI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = GOLD

    p = tf_tab.add_paragraph()
    p.text = "• Suami: 1/2 (tanpa anak) | 1/4 (bersama anak)\n• Istri: 1/4 (tanpa anak) | 1/8 (bersama anak)\n• Ayah: 'Ashabah penuh (tanpa anak) | 1/6 (ada anak)\n• Ibu: 1/3 (tanpa anak) | 1/6 (ada anak)\n• Anak Laki-laki & Perempuan: 2 : 1\n  (Li adz-dzakari mitslu hazhzhil untsayain)"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE

    c2 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = EMERALD

    tb_phi = s11.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.5))
    tf_phi = tb_phi.text_frame
    tf_phi.word_wrap = True
    p = tf_phi.paragraphs[0]
    p.text = "FALSAFAH KEADILAN PROPORSIONAL"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    p = tf_phi.add_paragraph()
    p.text = "1. Tanggung Jawab Finansial Pria:\nLaki-laki memikul kewajiban menafkahi istri, anak, tempat tinggal, mahar, dan keluarga besar.\n\n2. Hak Penuh Perempuan:\nWarisan perempuan adalah hak milik pribadinya secara utuh tanpa kewajiban menafkahi suami atau pihak lain.\n\n3. Banyak Kondisi Bagian Sama:\nAyah dan Ibu sama-sama mendapat 1/6 jika ada anak; saudara seibu laki-laki dan perempuan setara tanpa pembedaan rasio."
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    # ==========================================
    # SLIDE 12: KESIMPULAN
    # ==========================================
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s12)
    add_header(s12, "Kesimpulan Hasil Pembahasan Makalah", "BAB III • PENUTUP")

    sum_cards = [
        ("1. Hakikat Ahli Waris", "Ahli waris bukan sekadar famili biasa, melainkan pihak yang memenuhi syarat hukum syariat berdasarkan nasabiyah atau sababiyah.", GOLD),
        ("2. Klasifikasi Pria & Wanita", "Diakui 15 kelompok laki-laki & 10 perempuan. Ahli waris primer tak pernah gugur: Ayah, Ibu, Anak (Lk/Pr), & Suami/Istri.", EMERALD),
        ("3. Mekanisme & KHI", "Distribusi waris tunduk pada sistem Ashhabul Furudh, 'Ashabah, Hijab, serta selaras dengan Buku II Kompilasi Hukum Islam (KHI).", GOLD),
        ("4. Keadilan Transenden", "Syariat Islam menjamin hak ekonomi perempuan secara adil dan proporsional dengan beban tanggung jawab nafkah dalam peradaban Islam.", EMERALD),
    ]

    for i, (stitle, sdesc, col) in enumerate(sum_cards):
        x = Inches(0.8 + i * 2.95)
        c = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.8), Inches(5.0))
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = col

        tb = s12.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = LIGHT_GRAY

    # ==========================================
    # SLIDE 13: PENUTUP & DAFTAR PUSTAKA
    # ==========================================
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s13)

    c_end = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(6.5))
    c_end.fill.solid()
    c_end.fill.fore_color.rgb = CARD_BG
    c_end.line.color.rgb = GOLD
    c_end.line.width = Pt(2)

    tb_end = s13.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(10.9), Inches(5.8))
    tfe = tb_end.text_frame
    tfe.word_wrap = True

    p = tfe.paragraphs[0]
    p.text = "الْحَمْدُ لِلَّهِ رَبِّ الْعَالَمِينَ"
    p.font.size = Pt(26)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p = tfe.add_paragraph()
    p.text = "TERIMA KASIH & SESI TANYA JAWAB"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tfe.add_paragraph()
    p.text = "Disusun oleh Kelompok 5: Fatiaful Alzahra • Sufyan Tsaury • Wahdan Hamdun\nUIN Siber Syekh Nurjati Cirebon (2026)"
    p.font.size = Pt(14)
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER

    p = tfe.add_paragraph()
    p.text = "\nDaftar Pustaka Utama:\n• Al-Qur'an Surah An-Nisa (Ayat 7, 11, 12, 176)\n• Kitab Sahih Bukhari & Sahih Muslim (Bab Al-Fara'idh)\n• Inpres No. 1 Tahun 1991 tentang Kompilasi Hukum Islam (KHI) Buku II\n• Khair, A., & Zubair, A. (2022). Sistematika 'Asabah Dalam Hukum Kewarisan Islam\n• Widyaningrum, M. R., & Wiliyanarti, P. F. (2024). Kajian Kedudukan Perempuan dalam Kewarisan Islam"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Save
    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Ahli_Waris_Laki_Laki_dan_Perempuan.pptx")
    prs.save(out_path)
    print(f"File PPTX berhasil dibuat: {out_path}")

if __name__ == "__main__":
    create_presentation()
